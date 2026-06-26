from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

BG = "0F172A"
SURFACE = "111827"
PANEL = "1F2937"
PANEL_SOFT = "162032"
WHITE = "F0F0F0"
MUTED = "94A3B8"
LINE = "334155"
BLUE = "3B82F6"
BLUE_SOFT = "1D4ED8"
GREEN = "10B981"
AMBER = "F59E0B"
RED = "F43F5E"

FONT_CN = "Microsoft YaHei"
FONT_NUM = "Segoe UI"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, WIDE_W, WIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_glow_block(slide, left, top, width, height, color=BLUE_SOFT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = 0.82
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, size=20, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font=FONT_CN):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def style_paragraph(paragraph, size=18, color=WHITE, bold=False, font=FONT_CN, align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_kicker(slide, text, left=Inches(0.68), top=Inches(0.52), width=Inches(2.25), color=BLUE):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.36))
    pill.fill.solid()
    pill.fill.fore_color.rgb = rgb("12213C")
    pill.line.color.rgb = rgb(color)
    pill.line.width = Pt(1.25)
    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(10.5)
    run.font.bold = True
    run.font.color.rgb = rgb(color)
    return pill


def add_title_block(slide, title, subtitle, kicker, slide_no):
    add_kicker(slide, kicker)
    add_text(slide, title, Inches(0.68), Inches(1.0), Inches(7.4), Inches(0.95), size=28, bold=True)
    add_text(slide, subtitle, Inches(0.68), Inches(1.85), Inches(8.0), Inches(0.55), size=12.5, color=MUTED)
    no_chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(12.26), Inches(0.45), Inches(0.48), Inches(0.32))
    no_chip.fill.solid()
    no_chip.fill.fore_color.rgb = rgb("13233D")
    no_chip.line.color.rgb = rgb(LINE)
    tf = no_chip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{slide_no:02d}"
    run.font.name = FONT_NUM
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb(MUTED)


def add_footer(slide, source):
    add_text(slide, source, Inches(9.1), Inches(7.08), Inches(3.5), Inches(0.2), size=8.5, color="64748B", align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title=None, body=None, accent=None, fill=PANEL):
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.03), top + Inches(0.05), width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb("08101C")
    shadow.line.fill.background()

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(LINE)
    shape.line.width = Pt(1)

    if accent:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(accent)
        bar.line.fill.background()

    if title:
        add_text(slide, title, left + Inches(0.18), top + Inches(0.15), width - Inches(0.35), Inches(0.34), size=14.5, bold=True)
    if body:
        add_text(slide, body, left + Inches(0.18), top + Inches(0.52), width - Inches(0.35), height - Inches(0.65), size=10.8, color=MUTED)
    return shape


def add_metric(slide, left, top, width, height, value, label, color=BLUE):
    add_card(slide, left, top, width, height, accent=color, fill="182234")
    add_text(slide, value, left + Inches(0.16), top + Inches(0.18), width - Inches(0.3), Inches(0.32), size=20, bold=True, font=FONT_NUM)
    add_text(slide, label, left + Inches(0.16), top + Inches(0.52), width - Inches(0.3), Inches(0.22), size=10.2, color=MUTED)


def add_image_cover(slide, image_path, left, top, width, height):
    image_path = Path(image_path)
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h

    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.04), top + Inches(0.05), width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb("08101C")
    shadow.line.fill.background()

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = rgb("0B1220")
    frame.line.color.rgb = rgb("E5E7EB")
    frame.line.width = Pt(1)

    if img_ratio >= box_ratio:
        pic = slide.shapes.add_picture(str(image_path), left, top, height=height)
        overflow = pic.width - width
        if overflow > 0:
            pic.left = left - int(overflow / 2)
            crop = overflow / pic.width / 2
            pic.crop_left = crop
            pic.crop_right = crop
    else:
        pic = slide.shapes.add_picture(str(image_path), left, top, width=width)
        overflow = pic.height - height
        if overflow > 0:
            pic.top = top - int(overflow / 2)
            crop = overflow / pic.height / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
    return pic


def add_chip_row(slide, labels, left, top, width, gap=0.12, color=BLUE):
    chip_w = (width - Inches(gap) * (len(labels) - 1)) / len(labels)
    for idx, label in enumerate(labels):
        x = left + idx * (chip_w + Inches(gap))
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, chip_w, Inches(0.36))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb("13233D")
        chip.line.color.rgb = rgb(color)
        chip.line.width = Pt(1)
        tf = chip.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = FONT_CN
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = rgb(WHITE)


def add_arrow_line(slide, x1, y1, x2, y2, color=BLUE, width=2.25):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    try:
        conn.line.end_arrowhead = True
    except Exception:
        pass
    return conn


def build_deck():
    here = Path(__file__).resolve().parent
    repo_root = here.parents[3]
    screenshot_dir = repo_root / "screenshots"
    output_dir = here / "output"
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    slides_notes = []

    # 01 Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_glow_block(slide, Inches(7.75), Inches(0.78), Inches(4.7), Inches(5.25))
    add_kicker(slide, "中国大学生计算机设计大赛", width=Inches(1.75))
    add_text(slide, "学智画像\n把个性化学习真正做成一个闭环", Inches(0.72), Inches(1.2), Inches(5.6), Inches(1.75), size=30, bold=True)
    add_text(
        slide,
        "以学习画像为起点，把知识库 RAG、资源生成、学习规划与效果评估串成可演示、可回流、可持续优化的学习多智能体系统。",
        Inches(0.72), Inches(3.08), Inches(5.3), Inches(1.1), size=13, color=MUTED
    )
    add_image_cover(slide, screenshot_dir / "homepage.png", Inches(7.1), Inches(1.25), Inches(5.55), Inches(3.75))
    add_chip_row(slide, ["高校场景", "多智能体协同", "可信问答", "闭环学习"], Inches(0.72), Inches(4.48), Inches(5.45))
    add_metric(slide, Inches(0.72), Inches(5.2), Inches(1.42), Inches(0.88), "7类", "智能体协同", BLUE)
    add_metric(slide, Inches(2.25), Inches(5.2), Inches(1.42), Inches(0.88), "10维", "学习画像", GREEN)
    add_metric(slide, Inches(3.78), Inches(5.2), Inches(1.42), Inches(0.88), "8类", "资源生成", AMBER)
    add_metric(slide, Inches(5.31), Inches(5.2), Inches(1.42), Inches(0.88), "RAG", "可信问答", BLUE)
    add_footer(slide, "来源：README.md / screenshots/homepage.png")
    slides_notes.append(
        "第 1 页讲稿：这套系统的核心，不是单独做好某一个 AI 功能，而是把学习画像、资源生成、问答辅导、路径规划和评估反馈真正连成一条闭环。评委可以先记住一句话，学智画像想解决的不是“能不能生成”，而是“能不能持续支持一个学生学下去”。"
    )

    # 02 Pain
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "高校学习不是缺内容，而是缺少面向个人的连续支持", "问题并不在资源总量，而在于资源、诊断、计划和反馈之间长期是断开的。", "问题定义", 2)
    add_card(slide, Inches(0.7), Inches(2.3), Inches(3.85), Inches(2.45), "资源很多，但不一定适合我",
             "课程资料、题库、视频都不少，但学生很难快速找到和自己基础、目标、节奏真正匹配的内容。", BLUE)
    add_card(slide, Inches(4.74), Inches(2.3), Inches(3.85), Inches(2.45), "学习计划常常一刀切",
             "统一节奏很难兼顾学生差异，薄弱点、最佳学习时段和资源偏好通常没有被系统真正利用。", GREEN)
    add_card(slide, Inches(8.78), Inches(2.3), Inches(3.85), Inches(2.45), "评估结果难以回流到下一步",
             "很多系统把测试做成终点，只告诉学生分数，却没有把报告、资源推荐和后续计划接回去。", AMBER)
    add_card(slide, Inches(0.7), Inches(5.15), Inches(11.93), Inches(1.0), fill="132238")
    add_text(slide, "我们要解决的不是一个孤立功能，而是“诊断 - 生成 - 辅导 - 规划 - 再评估”这条支持链路为什么总是断开。", Inches(1.0), Inches(5.42), Inches(11.3), Inches(0.38), size=14, color=WHITE, bold=True)
    add_footer(slide, "来源：docs/A3_system_development_manual.md")
    slides_notes.append(
        "第 2 页讲稿：这一页建议直接从学生体验切入。学生并不是没有资源，而是不知道什么资源适合自己；教师也不是没有数据，而是很难把数据真正变成个性化动作。我们希望系统解决的，就是这种“看到了问题，却接不上下一步”的断裂。"
    )

    # 03 Loop
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "我们不是做聊天框，而是在做学习多智能体协同系统", "系统的价值不在单点回答，而在让不同智能体围绕同一个学生持续协作。", "总体方案", 3)
    stages = [
        ("学习画像", "理解目标、短板、偏好"),
        ("资源生成", "讲义 / PPT / 题库 / 视频"),
        ("问答辅导", "RAG 检索后再回答"),
        ("能力评估", "测评、解析、薄弱标签"),
        ("路径规划", "计划、提醒、资源推送"),
    ]
    x = Inches(0.86)
    for idx, (title, body) in enumerate(stages):
        add_card(slide, x + idx * Inches(2.45), Inches(2.65), Inches(2.0), Inches(1.3), title, body, BLUE if idx in (0, 2) else GREEN if idx == 4 else AMBER)
        if idx < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + idx * Inches(2.45) + Inches(2.04), Inches(3.05), Inches(0.28), Inches(0.42))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(BLUE)
            arrow.line.fill.background()
    add_card(slide, Inches(1.42), Inches(4.5), Inches(10.5), Inches(1.1), fill="132238")
    add_text(slide, "评估结果回流画像与计划，形成动态优化闭环，而不是一次性的静态输出。", Inches(1.7), Inches(4.83), Inches(9.8), Inches(0.24), size=13.5, bold=True)
    add_chip_row(slide, ["画像 Agent", "知识库 Agent", "生成 Agent", "规划 Agent"], Inches(0.82), Inches(6.0), Inches(7.1), gap=0.1, color=GREEN)
    add_chip_row(slide, ["辅导 Agent", "评估 Agent", "协调 Agent"], Inches(8.0), Inches(6.0), Inches(4.68), gap=0.1, color=GREEN)
    add_footer(slide, "来源：docs/A3_system_development_manual.md / docs/A3_ppt_video_outline.md")
    slides_notes.append(
        "第 3 页讲稿：这一页是全场最关键的一页。我们不是把几个 AI 功能简单堆在一起，而是让画像、资源、问答、评估和计划围绕同一个学生持续协同。这样评委就不会把项目理解成“一个做问答的页面”或者“一个做资源生成的页面”。"
    )

    # 04 Architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "四层架构让系统既能演示也能部署", "前端展示、业务路由、智能体协同和数据存储分层清晰，便于演示、扩展和落地部署。", "系统架构", 4)
    layers = [
        ("表现层", "首页 / 智能问答 / 能力测试 / 学习报告 / 学习计划"),
        ("业务路由层", "Flask + Blueprint + RESTful API，负责页面调度、任务分发与状态回传"),
        ("智能体层", "画像智能体 / 资源生成智能体 / 学习规划智能体 / 辅导智能体 / 评估智能体 / 知识库智能体 / 协调智能体"),
        ("数据层", "用户、答题记录、学习画像、学习资源、学习路径、知识库文档与分块索引"),
    ]
    y = Inches(2.18)
    fills = ["18243A", "162032", "132238", "18243A"]
    accents = [BLUE, GREEN, AMBER, BLUE]
    for idx, (title, body) in enumerate(layers):
        add_card(slide, Inches(0.88), y + idx * Inches(0.92), Inches(8.95), Inches(0.78), title, body, accents[idx], fill=fills[idx])
    add_card(slide, Inches(10.15), Inches(2.18), Inches(2.35), Inches(2.62), "部署能力",
             "Windows 一键启动\nDockerfile / docker-compose\nWSGI 生产部署", BLUE)
    add_card(slide, Inches(10.15), Inches(5.0), Inches(2.35), Inches(1.1), "技术栈",
             "Flask / SQLAlchemy / pandas / scikit-learn / ECharts / PyPDF2", GREEN)
    add_footer(slide, "来源：docs/A3_system_development_manual.md / README.md")
    slides_notes.append(
        "第 4 页讲稿：这一页建议用来回答“系统是不是只做了前端展示”。我们把前端页面、业务接口、智能体协同和数据存储分开设计，所以它既适合答辩演示，也具备继续扩展成稳定系统的工程基础。"
    )

    # 05 Profile
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "系统先理解学生，再决定给什么资源、怎么学", "画像不是静态标签，而是后续资源生成、问答策略和学习计划的共同输入。", "画像构建", 5)
    add_card(slide, Inches(0.82), Inches(2.25), Inches(3.3), Inches(3.65), "10 维画像如何服务后续决策", None, BLUE, fill="132238")
    dim_text = "专业背景 / 学习目标 / 薄弱知识点 / 兴趣方向 / 时间安排\n学习偏好 / 当前掌握度 / 风险点 / 资源偏好 / 进度节奏"
    add_text(slide, dim_text, Inches(1.03), Inches(2.8), Inches(2.85), Inches(1.45), size=14)
    add_card(slide, Inches(1.03), Inches(4.48), Inches(2.85), Inches(0.96), fill="18243A")
    add_text(slide, "对话输入 -> 特征抽取 -> 画像更新 -> 推荐触发", Inches(1.21), Inches(4.82), Inches(2.45), Inches(0.3), size=11.5, color=MUTED)
    add_image_cover(slide, screenshot_dir / "student_portrait.png", Inches(4.45), Inches(2.22), Inches(8.05), Inches(3.78))
    add_footer(slide, "来源：docs/A3_system_development_manual.md / screenshots/student_portrait.png")
    slides_notes.append(
        "第 5 页讲稿：先强调画像在这里不是一个展示面板，而是后续动作的起点。系统通过自然语言对话和学习数据共同更新画像，然后再决定应该生成什么资源、怎么答疑、学习计划怎么排，这样个性化才不是一句口号。"
    )

    # 06 Assistant hub
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "智能问答页其实是学习资源生成中枢", "在同一页面里把知识库问答、资源生成入口和多智能体协同操作集中起来，降低使用门槛。", "核心功能", 6)
    add_image_cover(slide, screenshot_dir / "intelligent_assistant_fixed.png", Inches(0.82), Inches(2.12), Inches(7.95), Inches(4.82))
    add_card(slide, Inches(9.05), Inches(2.2), Inches(3.05), Inches(1.16), "知识库问答", "优先基于课程资料检索，再生成回答，减少无依据输出。", BLUE)
    add_card(slide, Inches(9.05), Inches(3.52), Inches(3.05), Inches(1.5), "资源生成入口", "PPT、思维导图、课程讲解文档、题库、教学视频等资源可以从同一界面直接触发。", GREEN)
    add_card(slide, Inches(9.05), Inches(5.16), Inches(3.05), Inches(1.18), "多智能体协同", "画像、检索、生成、规划、评估在同一交互入口内完成联动。", AMBER)
    add_footer(slide, "来源：README.md / screenshots/intelligent_assistant_fixed.png")
    slides_notes.append(
        "第 6 页讲稿：这一页要帮助评委看到，智能问答不是一个孤立聊天框。它其实承担了统一入口的角色，学生从这里既可以问问题，也可以生成讲义、PPT、思维导图和题库，所以它更像学习多智能体系统的操作台。"
    )

    # 07 RAG
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "先检索再作答，让学习辅导更可信、可追溯", "系统把知识库上传、解析、分块、索引和引用回答接在一起，核心目标是降低幻觉风险。", "RAG 与可信回答", 7)
    process = ["上传资料", "解析分块", "索引检索", "引用回答", "结果回流"]
    base_x = Inches(0.92)
    for idx, label in enumerate(process):
        add_card(slide, base_x + idx * Inches(2.42), Inches(2.38), Inches(1.92), Inches(0.95), label, None, BLUE if idx in (0, 3) else GREEN if idx == 4 else AMBER, fill="18243A")
        if idx < len(process) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, base_x + idx * Inches(2.42) + Inches(1.96), Inches(2.64), Inches(0.3), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(BLUE)
            arrow.line.fill.background()
    add_card(slide, Inches(0.92), Inches(3.75), Inches(3.55), Inches(1.55), "支持的资料类型",
             "TXT / Markdown / CSV / JSON / DOCX / PDF，适合把课程讲义、实验文档和题库材料一起纳入知识库。", BLUE)
    add_card(slide, Inches(4.72), Inches(3.75), Inches(3.55), Inches(1.55), "可信回答策略",
             "回答优先依赖检索到的证据；未命中材料时，系统会提示资料不足，避免无依据生成。", GREEN)
    add_card(slide, Inches(8.52), Inches(3.75), Inches(3.55), Inches(1.55), "教学场景价值",
             "更适合高校课程辅导，因为学生看到的不只是答案，还有答案背后的依据和材料来源。", AMBER)
    add_image_cover(slide, screenshot_dir / "homepage.png", Inches(8.76), Inches(5.55), Inches(3.0), Inches(1.08))
    add_footer(slide, "来源：docs/A3_system_development_manual.md / docs/A3_document_checklist.md")
    slides_notes.append(
        "第 7 页讲稿：这里不要只说用了 RAG，而要说清楚它解决了什么问题。我们在教育场景下最在意的是可信度，所以系统先检索资料，再回答问题；如果资料不够，就提示材料不足，这样比直接生成更适合教学使用。"
    )

    # 08 Assessment
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "专业化测评让系统知道学生现在会什么、不会什么", "测评结果不只给出一个分数，而会拆成难度、方向、薄弱标签和逐题解析，方便后续继续干预。", "能力测试", 8)
    add_image_cover(slide, screenshot_dir / "assessment_fixed.png", Inches(0.82), Inches(2.18), Inches(7.75), Inches(4.95))
    add_card(slide, Inches(8.92), Inches(2.22), Inches(3.2), Inches(0.9), "10 类方向", "覆盖计算机、人工智能、数据分析及民航运行、空管、飞签派等方向。", BLUE)
    add_card(slide, Inches(8.92), Inches(3.28), Inches(3.2), Inches(0.9), "高校导向", "更贴近高校课程、专业能力与实验实训场景。", GREEN)
    add_card(slide, Inches(8.92), Inches(4.34), Inches(3.2), Inches(0.9), "逐题解析", "支持错题复盘、难度拆解与个性化建议。", AMBER)
    add_card(slide, Inches(8.92), Inches(5.4), Inches(3.2), Inches(0.9), "可回流画像", "测评结果继续回流学习报告、画像更新与学习计划重排。", BLUE)
    add_footer(slide, "来源：docs/A3_testing_spec.md / screenshots/assessment_fixed.png")
    slides_notes.append(
        "第 8 页讲稿：这一页建议把“专业化”三个字讲出来。我们不是做通用小测，而是把普通高校方向和民航特色方向一起纳入测评，输出逐题解析和薄弱标签，目的就是让系统真正知道学生下一步该补什么。"
    )

    # 09 Report and plan
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "测评结果不会停在报告里，而会继续回流到学习计划", "从结果洞察到任务安排，再到资源推荐和动态提醒，系统把“评估以后怎么办”这件事接住了。", "闭环回流", 9)
    add_card(slide, Inches(0.96), Inches(2.15), Inches(11.4), Inches(0.56), fill="132238")
    add_text(slide, "测评 -> 学习报告 -> 学习计划 -> 资源推荐 -> 再评估", Inches(2.55), Inches(2.33), Inches(8.2), Inches(0.2), size=12.5, color=WHITE, bold=True)
    add_image_cover(slide, screenshot_dir / "learning_report_fixed.png", Inches(0.82), Inches(2.95), Inches(5.82), Inches(3.62))
    add_image_cover(slide, screenshot_dir / "learning_plan_fixed.png", Inches(6.72), Inches(2.95), Inches(5.79), Inches(3.62))
    add_card(slide, Inches(0.82), Inches(6.72), Inches(3.78), Inches(0.5), "趋势洞察：让学生先看清状态变化，而不是只看一次成绩。 ", None, BLUE, fill="18243A")
    add_card(slide, Inches(4.79), Inches(6.72), Inches(3.78), Inches(0.5), "计划生成：根据掌握度、时间段和风险点动态调整任务节奏。", None, GREEN, fill="18243A")
    add_card(slide, Inches(8.76), Inches(6.72), Inches(3.78), Inches(0.5), "资源推送：把课程资源和学习动作一起推给学生，而不是分散跳转。", None, AMBER, fill="18243A")
    add_text(slide, "来源：A3_document_checklist.md / learning_report_fixed.png / learning_plan_fixed.png",
             Inches(8.1), Inches(1.92), Inches(4.1), Inches(0.2), size=8.0, color="64748B", align=PP_ALIGN.RIGHT)
    slides_notes.append(
        "第 9 页讲稿：很多系统把报告当成展示页，但我们的重点是报告之后怎么办。这里建议直接讲清楚，测评结果会回流到学习报告和学习计划，最后再带动资源推荐和后续复盘，这样才形成真正意义上的学习闭环。"
    )

    # 10 Readiness
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title_block(slide, "赛题要求与项目证据可以一页看清，答辩时可以一页讲透", "把能力点、完成度和证据入口集中在同一页面，帮助评委快速建立整体判断。", "赛题映射", 10)
    add_card(slide, Inches(0.82), Inches(2.2), Inches(5.15), Inches(4.9), fill="132238")
    rows = [
        ("对话式学习画像", "已完成", GREEN),
        ("多智能体协同资源生成", "已完成", GREEN),
        ("个性化路径规划与资源推送", "已完成", GREEN),
        ("RAG 智能辅导", "已完成", GREEN),
        ("防幻觉与课程样例包", "进行中", AMBER),
    ]
    add_text(slide, "关键赛题要求映射", Inches(1.06), Inches(2.48), Inches(2.5), Inches(0.3), size=15, bold=True)
    y = Inches(3.0)
    for title, status, color in rows:
        add_text(slide, title, Inches(1.06), y, Inches(2.95), Inches(0.24), size=12.3, bold=True)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.15), y - Inches(0.02), Inches(0.92), Inches(0.28))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb("14263E")
        chip.line.color.rgb = rgb(color)
        tf = chip.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = status
        run.font.name = FONT_CN
        run.font.size = Pt(9.5)
        run.font.bold = True
        run.font.color.rgb = rgb(color)
        y += Inches(0.7)
    add_image_cover(slide, screenshot_dir / "homepage.png", Inches(6.28), Inches(2.2), Inches(6.0), Inches(4.9))
    add_footer(slide, "来源：screenshots/homepage.png")
    slides_notes.append(
        "第 10 页讲稿：这一页适合答辩时快速收束。评委往往会关心项目到底和赛题贴不贴，这里可以直接用页面和要点一一对应，让大家看到我们不仅功能上契合，而且已经把证据整理到适合答辩说明的形式。"
    )

    # 11 Roadmap
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "阶段计划")
    add_text(slide, "当前核心闭环已成型，下一步重点补齐提交物与证据包", Inches(0.68), Inches(0.98), Inches(8.55), Inches(1.25), size=25, bold=True)
    add_text(slide, "后续工作不再是大改系统主线，而是把样例包、答辩材料和证据链补齐，让提交更完整。", Inches(0.68), Inches(2.0), Inches(8.1), Inches(0.34), size=12.2, color=MUTED)
    no_chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(12.26), Inches(0.45), Inches(0.48), Inches(0.32))
    no_chip.fill.solid()
    no_chip.fill.fore_color.rgb = rgb("13233D")
    no_chip.line.color.rgb = rgb(LINE)
    tf = no_chip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "11"
    run.font.name = FONT_NUM
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb(MUTED)
    lanes = [
        ("P0 必须优先完成", [
            "补一门完整专业课程的知识库样例包",
            "完成演示 PPT 终稿与 7 分钟脚本",
            "录制并剪辑 7 分钟内演示视频",
            "补测试截图、流程截图和接口验证截图",
        ], BLUE),
        ("P1 建议继续增强", [
            "补“防幻觉与内容安全过滤机制”说明",
            "补资源生成模块的实际案例截图",
            "补能力测试、学习计划、学习报告使用路径图",
        ], GREEN),
        ("P2 提交前整理", [
            "Markdown 草稿转 Word / PDF 正式版",
            "统一封面、编号、图表标题和术语表",
            "按赛题要求做最终逐项核查",
        ], AMBER),
    ]
    top = Inches(2.35)
    for idx, (title, items, color) in enumerate(lanes):
        y = top + idx * Inches(1.55)
        add_card(slide, Inches(0.84), y, Inches(12.0), Inches(1.25), title, None, color, fill="132238")
        for item_idx, item in enumerate(items):
            add_text(slide, f"{item_idx + 1}. {item}", Inches(1.1), y + Inches(0.38) + item_idx * Inches(0.24), Inches(10.0), Inches(0.2), size=10.6, color=WHITE)
    add_card(slide, Inches(9.0), Inches(5.92), Inches(3.2), Inches(0.98), "答辩落点", "讲清闭环、展示系统、证明契合，再用真实页面和材料建立可信度。", BLUE)
    add_footer(slide, "来源：docs/A3_document_checklist.md")
    slides_notes.append(
        "第 11 页讲稿：最后一页不要给人“项目还没做完”的感觉，而是要表达“核心闭环已经完成，现在正在补齐提交物和证据链”。这样既诚实，也能让评委相信项目已经具备完整的演示基础。"
    )

    out_path = output_dir / "学智画像_A3答辩初稿_深墨蓝科技风.pptx"
    prs.save(out_path)

    notes_lines = ["# 学智画像 A3 答辩初稿讲稿", ""]
    for idx, note in enumerate(slides_notes, start=1):
        notes_lines.append(f"## 第 {idx} 页")
        notes_lines.append(note)
        notes_lines.append("")
    (output_dir / "学智画像_A3答辩初稿_讲稿备注.md").write_text("\n".join(notes_lines), encoding="utf-8")

    print(out_path)


if __name__ == "__main__":
    build_deck()
